# OPÇÃO DE EBOOK DESABILITADA ATE SEGUNDA ORDEM
import openpyxl
from pptx import Presentation
from ebooklib import epub
from docx import Document
from docx2pdf import convert
import os

# Função para extrair palavras-chave de planilhas Excel
def extract_keywords_from_excel(file_path, sheet_name='Semana4', column=1):
    keywords = []
    wb = openpyxl.load_workbook(file_path)
    sheet = wb[sheet_name]

    for row in sheet.iter_rows(min_col=3, max_col=3, values_only=True):
        if row[0]:
            keywords.append(row[0])  # Adiciona as palavras-chave de cada linha

    return keywords

# Função para extrair palavras-chave de apresentações PowerPoint
def extract_keywords_from_pptx(pptx_path):
    keywords = []
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                keywords.extend(shape.text.split())  # Extrai as palavras-chave de cada slide
    return keywords

# Função para extrair palavras-chave de eBooks (ePub)
# def extract_keywords_from_epub(epub_path):
#     keywords = []
#     book = epub.read_epub(epub_path)
#     for item in book.get_items():
#         if item.get_type() == epub.ITEM_DOCUMENT:
#             text = item.get_content().decode('utf-8')
#             keywords.extend(text.split())  # Extrai palavras-chave do eBook
#     return keywords

# Função para criar um documento Word com as palavras-chave
def create_doc_with_keywords(keywords, output_path):
    doc = Document()
    doc.add_heading('Palavras-Chave', level=1)

    for keyword in keywords:
        doc.add_paragraph(keyword)

    doc.save(output_path)
    print(f"Arquivo .docx criado: {output_path}")

# Função para converter o documento .docx para PDF
def convert_doc_to_pdf(input_path, output_path=None):
    if output_path is None:
        output_path = os.path.splitext(input_path)[0] + ".pdf"

    convert(input_path, output_path)
    print(f"Arquivo PDF criado: {output_path}")

# Exemplo de uso
def main():
    # Caminhos para os arquivos
    excel_file = 'ODS - Entregas.xlsx'  # Planilha com respostas do formulário
    pptx_file = 'projeto teste.pptx'           # Apresentação PowerPoint
    # epub_file = 'ebook.epub'                  # eBook no formato .epub
    docx_file = 'palavras_chaves.docx'        # Documento Word de saída

    # 1. Extrair palavras-chave de cada tipo de arquivo
    excel_keywords = extract_keywords_from_excel(excel_file)
    pptx_keywords = extract_keywords_from_pptx(pptx_file)
    # epub_keywords = extract_keywords_from_epub(epub_file)

    # 2. Combinar todas as palavras-chave
    all_keywords = excel_keywords + pptx_keywords
    # + epub_keywords

    # 3. Criar o arquivo .docx com as palavras-chave
    create_doc_with_keywords(all_keywords, docx_file)

    # 4. Converter o arquivo .docx para PDF
    convert_doc_to_pdf(docx_file)

if __name__ == "__main__":
    main()
